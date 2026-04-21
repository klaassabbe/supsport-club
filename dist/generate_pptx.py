"""
SupSport Club — Générateur de présentation PowerPoint
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.util as util

# ── Couleurs ──────────────────────────────────────────────────────
TEAL       = RGBColor(0x00, 0x9E, 0xBE)  # #009EBE
DARK_TEAL  = RGBColor(0x00, 0x68, 0x80)  # #006880
NAVY       = RGBColor(0x00, 0x3D, 0x50)  # #003d50
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF1, 0xF5, 0xF9)
MID_GRAY   = RGBColor(0x64, 0x74, 0x8B)
DARK       = RGBColor(0x1A, 0x1A, 0x2E)
GREEN      = RGBColor(0x16, 0xA3, 0x4A)
YELLOW     = RGBColor(0xF5, 0xBF, 0x00)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

# ── Helper functions ──────────────────────────────────────────────
def blank(prs):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)

def bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def gradient_bg(slide, c1: RGBColor, c2: RGBColor):
    """Approximated with solid (pptx doesn't support true gradient bg easily)"""
    bg(slide, c1)

def rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, left, top, width, height,
        size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_bullet_box(slide, items, left, top, width, height,
                   check_color=TEAL, text_color=DARK, size=13, check="✓"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        # check mark
        r1 = p.add_run()
        r1.text = check + "  "
        r1.font.size = Pt(size)
        r1.font.bold = True
        r1.font.color.rgb = check_color
        # text
        r2 = p.add_run()
        r2.text = item
        r2.font.size = Pt(size)
        r2.font.color.rgb = text_color

# ═══════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ═══════════════════════════════════════════════════════════════════
sl = blank(prs)
bg(sl, NAVY)

# Accent rectangle left
rect(sl, 0, 0, Inches(0.35), SLIDE_H, fill_color=TEAL)

# Decorative circle top-right
sh = sl.shapes.add_shape(9, Inches(9.5), Inches(-1.5), Inches(5), Inches(5))
sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(0x00, 0x68, 0x80)
sh.line.fill.background()

# Badge
r = rect(sl, Inches(0.8), Inches(1.2), Inches(4.5), Inches(0.4),
         fill_color=RGBColor(0x00, 0x5E, 0x72))
t = txt(sl, "PROPOSITION DE PARTENARIAT — CONFIDENTIEL",
        Inches(0.8), Inches(1.18), Inches(4.5), Inches(0.45),
        size=9, bold=True, color=RGBColor(0x69, 0xC3, 0xD2), align=PP_ALIGN.CENTER)

# Main title
txt(sl, "Financez vos athlètes.", Inches(0.8), Inches(1.9), Inches(9), Inches(1.1),
    size=48, bold=True, color=WHITE)
txt(sl, "Sans budget.", Inches(0.8), Inches(2.9), Inches(9), Inches(1.1),
    size=48, bold=True, color=RGBColor(0x69, 0xC3, 0xD2))

# Subtitle
txt(sl, "SupSport Club transforme la communauté de votre fédération en moteur\nde financement — et vous apporte visibilité, communication et revenus récurrents.",
    Inches(0.8), Inches(4.0), Inches(8.5), Inches(1.2),
    size=16, color=RGBColor(0xB0, 0xC8, 0xD4))

# 4 key numbers
stats = [("0 €", "investissement initial"), ("10%", "reversé à SupSport"), ("13", "partenaires exclusifs"), ("100%", "gagnant-gagnant")]
for i, (val, lbl) in enumerate(stats):
    x = Inches(0.8 + i * 3.1)
    rect(sl, x, Inches(5.5), Inches(2.8), Inches(1.6),
         fill_color=RGBColor(0x00, 0x55, 0x6A))
    txt(sl, val,  x, Inches(5.55), Inches(2.8), Inches(0.8),
        size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, lbl,  x, Inches(6.2), Inches(2.8), Inches(0.7),
        size=11, color=RGBColor(0x69, 0xC3, 0xD2), align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 2 — LE PROBLÈME
# ═══════════════════════════════════════════════════════════════════
sl = blank(prs)
bg(sl, LIGHT_GRAY)

rect(sl, 0, 0, SLIDE_W, Inches(0.06), fill_color=TEAL)

txt(sl, "LE DÉFI DE TOUTES LES FÉDÉRATIONS", Inches(0.7), Inches(0.3), Inches(8), Inches(0.4),
    size=10, bold=True, color=TEAL)
txt(sl, "Un problème que vous connaissez bien", Inches(0.7), Inches(0.65), Inches(11), Inches(0.8),
    size=30, bold=True, color=DARK)

# 3 problem cards
problems = [
    ("💸", "Budget insuffisant", "Vos athlètes abandonnent faute de moyens pour financer inscriptions, équipements et déplacements."),
    ("👁️", "Manque de visibilité", "Votre fédération peine à se faire connaître du grand public et à fidéliser ses membres existants."),
    ("🤝", "Sponsors difficiles", "Les partenariats traditionnels prennent des mois à négocier et les montants restent incertains."),
]
for i, (icon, title, desc) in enumerate(problems):
    x = Inches(0.5 + i * 4.2)
    r = rect(sl, x, Inches(1.7), Inches(3.9), Inches(2.8),
             fill_color=WHITE, line_color=RGBColor(0xE2, 0xE8, 0xF0))
    txt(sl, icon, x + Inches(0.2), Inches(1.85), Inches(0.8), Inches(0.7), size=28, color=DARK)
    txt(sl, title, x + Inches(0.2), Inches(2.5), Inches(3.5), Inches(0.5),
        size=15, bold=True, color=DARK)
    txt(sl, desc, x + Inches(0.2), Inches(2.95), Inches(3.5), Inches(1.3),
        size=12, color=MID_GRAY)

# Solution intro
rect(sl, Inches(0.5), Inches(4.7), SLIDE_W - Inches(1), Inches(2.3),
     fill_color=NAVY)
txt(sl, "La solution SupSport Club :", Inches(0.9), Inches(4.85), Inches(10), Inches(0.5),
    size=14, bold=True, color=RGBColor(0x69, 0xC3, 0xD2))
txt(sl, "Un programme de financement participatif où vos supporters paient un abonnement mensuel\nen échange d'avantages exclusifs chez 13 partenaires sportifs — et vos athlètes reçoivent directement les fonds.",
    Inches(0.9), Inches(5.3), Inches(11.5), Inches(1.4),
    size=14, color=WHITE)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 3 — WIN WIN
# ═══════════════════════════════════════════════════════════════════
sl = blank(prs)
bg(sl, WHITE)
rect(sl, 0, 0, SLIDE_W, Inches(0.06), fill_color=TEAL)

txt(sl, "MODÈLE GAGNANT-GAGNANT", Inches(0.7), Inches(0.3), Inches(8), Inches(0.4),
    size=10, bold=True, color=TEAL)
txt(sl, "Tout le monde y gagne — vraiment.", Inches(0.7), Inches(0.65), Inches(11), Inches(0.8),
    size=30, bold=True, color=DARK)

# Left card — Fédération
rect(sl, Inches(0.5), Inches(1.6), Inches(5.9), Inches(5.4),
     fill_color=RGBColor(0xF0, 0xFB, 0xFF), line_color=RGBColor(0x00, 0x9E, 0xBE),
     line_width=Pt(1.5))
txt(sl, "🏆  Pour la fédération", Inches(0.7), Inches(1.75), Inches(5.5), Inches(0.5),
    size=15, bold=True, color=TEAL)
fed_items = [
    "Athlètes financés sans toucher au budget fédération",
    "Site professionnel à vos couleurs en quelques jours",
    "Jusqu'à 5% de chaque don reversé à la fédération",
    "Communication régulière : affiches, newsletter (selon formule)",
    "Image moderne et innovante face aux autres fédérations",
    "Tableau de bord : dons, athlètes, supporters en temps réel",
]
add_bullet_box(sl, fed_items, Inches(0.7), Inches(2.3), Inches(5.5), Inches(4.2),
               check_color=TEAL, text_color=DARK, size=13)

# Right card — Supporters
rect(sl, Inches(6.7), Inches(1.6), Inches(5.9), Inches(5.4),
     fill_color=RGBColor(0xF0, 0xFD, 0xF4), line_color=GREEN,
     line_width=Pt(1.5))
txt(sl, "❤️  Pour les supporters", Inches(6.9), Inches(1.75), Inches(5.5), Inches(0.5),
    size=15, bold=True, color=GREEN)
supp_items = [
    "Soutien concret à des athlètes locaux qu'ils connaissent",
    "Accès à 13 partenaires sportifs (remises 10% à 20%)",
    "Le don peut être rentabilisé dès le 1er mois",
    "Sentiment d'appartenance à une vraie communauté",
    "Contenus exclusifs et actualités des athlètes",
    "Badge officiel de supporter de la fédération",
]
add_bullet_box(sl, supp_items, Inches(6.9), Inches(2.3), Inches(5.5), Inches(4.2),
               check_color=GREEN, text_color=DARK, size=13)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 4 — COMMENT ÇA FONCTIONNE
# ═══════════════════════════════════════════════════════════════════
sl = blank(prs)
bg(sl, LIGHT_GRAY)
rect(sl, 0, 0, SLIDE_W, Inches(0.06), fill_color=TEAL)

txt(sl, "FONCTIONNEMENT", Inches(0.7), Inches(0.3), Inches(8), Inches(0.4),
    size=10, bold=True, color=TEAL)
txt(sl, "Simple. Automatique. Efficace.", Inches(0.7), Inches(0.65), Inches(11), Inches(0.8),
    size=30, bold=True, color=DARK)

steps = [
    ("🌐", "01", "La page de votre club", "Un site sur mesure aux couleurs de votre fédération avec les profils de tous vos athlètes, leurs objectifs et leurs palmarès."),
    ("💳", "02", "Le supporter choisit son niveau", "De 10€ à 25€/mois (ou don unique). Il choisit ses avantages partenaires. Son don est automatiquement réparti entre les athlètes."),
    ("🏅", "03", "Tout le monde gagne", "L'athlète reçoit son financement. Le supporter économise sur ses achats. La fédération perçoit 5% si elle l'active."),
]
for i, (icon, num, title, desc) in enumerate(steps):
    x = Inches(0.5 + i * 4.2)
    rect(sl, x, Inches(1.6), Inches(3.9), Inches(3.8), fill_color=WHITE,
         line_color=RGBColor(0xE2, 0xE8, 0xF0))
    # Step number (big watermark)
    txt(sl, num, x + Inches(2.5), Inches(1.6), Inches(1.2), Inches(1.2),
        size=48, bold=True, color=RGBColor(0xE2, 0xE8, 0xF0))
    txt(sl, icon, x + Inches(0.2), Inches(1.75), Inches(0.8), Inches(0.7), size=26, color=DARK)
    txt(sl, title, x + Inches(0.2), Inches(2.45), Inches(3.5), Inches(0.6),
        size=14, bold=True, color=DARK)
    txt(sl, desc, x + Inches(0.2), Inches(3.05), Inches(3.5), Inches(2.0),
        size=12, color=MID_GRAY)

# Arrow connectors (simple text)
txt(sl, "→", Inches(4.45), Inches(3.1), Inches(0.5), Inches(0.5), size=24, bold=True, color=TEAL)
txt(sl, "→", Inches(8.65), Inches(3.1), Inches(0.5), Inches(0.5), size=24, bold=True, color=TEAL)

# Répartition du don
rect(sl, Inches(0.5), Inches(5.7), SLIDE_W - Inches(1), Inches(1.45), fill_color=NAVY)
txt(sl, "💡  Répartition de chaque don :", Inches(0.8), Inches(5.8), Inches(5), Inches(0.45),
    size=13, bold=True, color=RGBColor(0x69, 0xC3, 0xD2))
splits = [("~89%", "aux athlètes", GREEN), ("5%", "au club (optionnel)", TEAL), ("~10%", "plateforme SupSport", YELLOW)]
for i, (pct, lbl, col) in enumerate(splits):
    x = Inches(3.5 + i * 3.2)
    txt(sl, pct, x, Inches(5.82), Inches(1.5), Inches(0.55), size=20, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(sl, lbl, x, Inches(6.3), Inches(1.5), Inches(0.4), size=10, color=WHITE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 5 — SIMULATION FINANCIÈRE
# ═══════════════════════════════════════════════════════════════════
sl = blank(prs)
bg(sl, NAVY)
rect(sl, 0, 0, SLIDE_W, Inches(0.06), fill_color=TEAL)

txt(sl, "SIMULATION FINANCIÈRE", Inches(0.7), Inches(0.3), Inches(8), Inches(0.4),
    size=10, bold=True, color=RGBColor(0x69, 0xC3, 0xD2))
txt(sl, "Ce que ça représente concrètement", Inches(0.7), Inches(0.65), Inches(11), Inches(0.8),
    size=30, bold=True, color=WHITE)
txt(sl, "Exemple réaliste : 60 supporters actifs au niveau Champion (€25/mois)",
    Inches(0.7), Inches(1.35), Inches(11), Inches(0.45), size=14, color=RGBColor(0xB0, 0xC8, 0xD4))

# Left col — Fédération
rect(sl, Inches(0.5), Inches(2.0), Inches(5.9), Inches(4.8),
     fill_color=RGBColor(0x00, 0x4D, 0x63))
txt(sl, "Pour la fédération & les athlètes", Inches(0.7), Inches(2.1), Inches(5.5), Inches(0.5),
    size=13, bold=True, color=RGBColor(0x69, 0xC3, 0xD2))

rows_fed = [
    ("Dons bruts collectés / mois",   "€1 500",   WHITE),
    ("Part reversée aux athlètes (~89%)", "€1 335", GREEN),
    ("Part reversée à la fédération (5%)", "€75",   GREEN),
    ("Soit en 12 mois pour la fédération", "€900",  YELLOW),
    ("+ Valeur comm. & visibilité",    "inestimable", RGBColor(0x69,0xC3,0xD2)),
]
for j, (lbl, val, col) in enumerate(rows_fed):
    y = Inches(2.7 + j * 0.72)
    rect(sl, Inches(0.5), y, Inches(5.9), Inches(0.65),
         fill_color=RGBColor(0x00, 0x56, 0x6E) if j % 2 == 0 else RGBColor(0x00, 0x4D, 0x63))
    txt(sl, lbl, Inches(0.7), y + Inches(0.12), Inches(3.8), Inches(0.45), size=12, color=WHITE)
    txt(sl, val, Inches(4.3), y + Inches(0.12), Inches(2.0), Inches(0.45), size=13, bold=True,
        color=col, align=PP_ALIGN.RIGHT)

# Right col — Supporter
rect(sl, Inches(6.7), Inches(2.0), Inches(5.9), Inches(4.8),
     fill_color=RGBColor(0x00, 0x4D, 0x63))
txt(sl, "Pour chaque supporter (niveau Champion)", Inches(6.9), Inches(2.1), Inches(5.5), Inches(0.5),
    size=13, bold=True, color=RGBColor(0x69, 0xC3, 0xD2))

rows_supp = [
    ("Don mensuel",                         "-€25",        RGBColor(0xB0,0xC8,0xD4)),
    ("Remise XRUN chaussures running",      "+€36",        GREEN),
    ("Remise Nutri-bay commande mensuelle", "+€8",         GREEN),
    ("Remise Physiosport (séance récup)",   "+€7,50",      GREEN),
]
for j, (lbl, val, col) in enumerate(rows_supp):
    y = Inches(2.7 + j * 0.72)
    rect(sl, Inches(6.7), y, Inches(5.9), Inches(0.65),
         fill_color=RGBColor(0x00, 0x56, 0x6E) if j % 2 == 0 else RGBColor(0x00, 0x4D, 0x63))
    txt(sl, lbl, Inches(6.9), y + Inches(0.12), Inches(3.8), Inches(0.45), size=12, color=WHITE)
    txt(sl, val, Inches(10.3), y + Inches(0.12), Inches(2.0), Inches(0.45), size=13, bold=True,
        color=col, align=PP_ALIGN.RIGHT)

# Total supporter
rect(sl, Inches(6.7), Inches(5.6), Inches(5.9), Inches(0.9), fill_color=GREEN)
txt(sl, "Économie nette par mois", Inches(6.9), Inches(5.7), Inches(3.5), Inches(0.6),
    size=14, bold=True, color=WHITE)
txt(sl, "+€26,50 🎉", Inches(9.8), Inches(5.65), Inches(2.5), Inches(0.7),
    size=18, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 6 — LES 3 FORMULES (titre)
# ═══════════════════════════════════════════════════════════════════
sl = blank(prs)
bg(sl, WHITE)
rect(sl, 0, 0, SLIDE_W, Inches(0.06), fill_color=TEAL)

txt(sl, "NOS FORMULES", Inches(0.7), Inches(0.3), Inches(8), Inches(0.4),
    size=10, bold=True, color=TEAL)
txt(sl, "Choisissez votre niveau d'accompagnement", Inches(0.7), Inches(0.65), Inches(11), Inches(0.8),
    size=30, bold=True, color=DARK)
txt(sl, "Sans engagement longue durée · Préavis d'un mois",
    Inches(0.7), Inches(1.35), Inches(8), Inches(0.4), size=13, color=MID_GRAY)

plans = [
    {
        "icon": "🌐", "name": "ESSENTIEL", "price": "€250", "sub": "/mois",
        "desc": "Votre fédération sur le web avec un système de don clé en main.",
        "color": RGBColor(0xE2,0xE8,0xF0), "text": DARK, "featured": False,
        "features": [
            "Site personnalisé aux couleurs du club",
            "Profils illimités d'athlètes",
            "Système de don intégré",
            "Accès aux 13 partenaires SupSport",
            "Option 5% pour le club activable",
            "Tableau de bord en temps réel",
            "Support par email",
        ]
    },
    {
        "icon": "📣", "name": "PRO ⭐", "price": "€550", "sub": "/mois",
        "desc": "Site + présence physique et digitale pour faire connaître votre programme.",
        "color": NAVY, "text": WHITE, "featured": True,
        "features": [
            "Tout ce qui est dans Essentiel",
            "Visite trimestrielle dans votre fédération",
            "1 affiche professionnelle / trimestre",
            "Newsletter mensuelle pour vos membres",
            "Rapport mensuel : dons, supporters, ROI",
            "Mise à jour des profils athlètes incluse",
            "Support prioritaire (réponse 24h)",
        ]
    },
    {
        "icon": "🚀", "name": "PREMIUM", "price": "€1 000", "sub": "/mois",
        "desc": "Accompagnement complet. Nous devenons votre partenaire communication.",
        "color": RGBColor(0xF8, 0xFA, 0xFC), "text": DARK, "featured": False,
        "features": [
            "Tout ce qui est dans Pro",
            "Affiche mensuelle (12/an)",
            "Campagne réseaux sociaux mensuelle",
            "Visite mensuelle dans votre fédération",
            "Extension aux clubs affiliés",
            "Campagne recrutement de supporters",
            "Account manager dédié",
        ]
    },
]

for i, plan in enumerate(plans):
    x = Inches(0.3 + i * 4.35)
    w = Inches(4.1)
    card_color = plan["color"]
    is_featured = plan["featured"]

    # Card background
    rect(sl, x, Inches(1.9), w, Inches(5.2),
         fill_color=card_color,
         line_color=TEAL if is_featured else RGBColor(0xE2,0xE8,0xF0),
         line_width=Pt(2) if is_featured else Pt(1))

    if is_featured:
        # "Recommandé" badge
        rect(sl, x + Inches(1.1), Inches(1.65), Inches(1.9), Inches(0.32), fill_color=TEAL)
        txt(sl, "⭐ RECOMMANDÉ", x + Inches(1.1), Inches(1.66), Inches(1.9), Inches(0.3),
            size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    tc = plan["text"]
    txt(sl, plan["icon"], x + Inches(0.2), Inches(2.0), Inches(0.7), Inches(0.6), size=22, color=tc)
    txt(sl, plan["name"], x + Inches(0.2), Inches(2.55), Inches(3.7), Inches(0.4),
        size=11, bold=True, color=TEAL if not is_featured else RGBColor(0x69,0xC3,0xD2))
    txt(sl, plan["price"], x + Inches(0.2), Inches(2.9), Inches(2.5), Inches(0.65),
        size=32, bold=True, color=tc)
    txt(sl, plan["sub"], x + Inches(0.2), Inches(3.45), Inches(1.5), Inches(0.35),
        size=12, color=MID_GRAY if not is_featured else RGBColor(0xB0,0xC8,0xD4))
    txt(sl, plan["desc"], x + Inches(0.2), Inches(3.8), Inches(3.7), Inches(0.7),
        size=11, color=MID_GRAY if not is_featured else RGBColor(0xB0,0xC8,0xD4))

    add_bullet_box(sl, plan["features"],
                   x + Inches(0.15), Inches(4.55), Inches(3.8), Inches(2.3),
                   check_color=TEAL, text_color=tc, size=11)

# Commission note
rect(sl, Inches(0.3), Inches(7.1), SLIDE_W - Inches(0.6), Inches(0.35),
     fill_color=RGBColor(0xF0,0xFB,0xFF), line_color=RGBColor(0x69,0xC3,0xD2))
txt(sl, "⚡  Sur chaque don collecté, ~10% revient à SupSport pour couvrir les frais de plateforme (Stripe, hébergement, développement) — ce modèle finance l'amélioration continue de votre site.",
    Inches(0.5), Inches(7.1), SLIDE_W - Inches(1), Inches(0.35),
    size=10, color=TEAL, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 7 — OBJECTIONS / FAQ
# ═══════════════════════════════════════════════════════════════════
sl = blank(prs)
bg(sl, LIGHT_GRAY)
rect(sl, 0, 0, SLIDE_W, Inches(0.06), fill_color=TEAL)

txt(sl, "QUESTIONS FRÉQUENTES", Inches(0.7), Inches(0.3), Inches(8), Inches(0.4),
    size=10, bold=True, color=TEAL)
txt(sl, "Vous vous posez peut-être ces questions", Inches(0.7), Inches(0.65), Inches(11), Inches(0.8),
    size=30, bold=True, color=DARK)

faqs = [
    ("Est-ce que nos membres vont vraiment souscrire ?",
     "Notre modèle est conçu pour que le don soit rentable pour le supporter : les remises partenaires lui rapportent souvent plus que ce qu'il donne. Ce n'est pas de la charité, c'est une offre d'avantages."),
    ("Combien de temps pour mettre en place ?",
     "5 à 10 jours ouvrés après réception des infos du club et profils athlètes. Pas de compétences techniques nécessaires de votre côté — on s'occupe de tout."),
    ("Peut-on arrêter quand on veut ?",
     "Oui. Pas d'engagement annuel imposé. Un préavis d'un mois suffit. Nous préférons que vous restiez parce que ça fonctionne, pas parce que vous y êtes obligés."),
    ("Est-ce que les 250€/mois s'équilibrent avec la commission ?",
     "Avec 60 supporters à 25€/mois, vous collectez 1 500€. La plateforme prend 150€, le club récupère 75€ (option 5%). Vos athlètes reçoivent 1 275€ de plus que sans SupSport."),
]

for i, (q, a) in enumerate(faqs):
    row = i // 2
    col = i % 2
    x = Inches(0.5 + col * 6.4)
    y = Inches(1.7 + row * 2.6)
    rect(sl, x, y, Inches(6.0), Inches(2.35), fill_color=WHITE,
         line_color=RGBColor(0xE2,0xE8,0xF0))
    # Q badge
    rect(sl, x + Inches(0.2), y + Inches(0.18), Inches(0.28), Inches(0.28), fill_color=TEAL)
    txt(sl, "Q", x + Inches(0.2), y + Inches(0.17), Inches(0.28), Inches(0.28),
        size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, q, x + Inches(0.6), y + Inches(0.15), Inches(5.2), Inches(0.55),
        size=13, bold=True, color=DARK)
    txt(sl, a, x + Inches(0.2), y + Inches(0.72), Inches(5.6), Inches(1.5),
        size=12, color=MID_GRAY)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 8 — CTA FINAL
# ═══════════════════════════════════════════════════════════════════
sl = blank(prs)
bg(sl, NAVY)

# Decorative circle
sh = sl.shapes.add_shape(9, Inches(8.5), Inches(-2), Inches(6), Inches(6))
sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(0x00, 0x55, 0x6A)
sh.line.fill.background()

rect(sl, 0, 0, SLIDE_W, Inches(0.06), fill_color=TEAL)

# Central content
txt(sl, "Prêt à lancer le programme ?", Inches(1), Inches(1.8), Inches(11), Inches(1.2),
    size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(sl, "La démo est déjà en ligne. Vos athlètes méritent ce soutien.",
    Inches(2), Inches(3.0), Inches(9), Inches(0.6),
    size=17, color=RGBColor(0xB0, 0xC8, 0xD4), align=PP_ALIGN.CENTER)

# Two CTA boxes
rect(sl, Inches(2.2), Inches(3.9), Inches(3.8), Inches(0.9), fill_color=TEAL)
txt(sl, "🌐  Voir la démo en ligne", Inches(2.2), Inches(3.95), Inches(3.8), Inches(0.75),
    size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rect(sl, Inches(7.3), Inches(3.9), Inches(3.8), Inches(0.9),
     fill_color=RGBColor(0x00, 0x00, 0x00))
rect(sl, Inches(7.3), Inches(3.9), Inches(3.8), Inches(0.9),
     fill_color=RGBColor(0x00, 0x55, 0x6A), line_color=WHITE, line_width=Pt(1))
txt(sl, "✉️  klaas@supsport.be", Inches(7.3), Inches(3.95), Inches(3.8), Inches(0.75),
    size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Next steps
rect(sl, Inches(1.5), Inches(5.1), Inches(10.3), Inches(2.0),
     fill_color=RGBColor(0x00, 0x4D, 0x63))
txt(sl, "Prochaines étapes :", Inches(1.8), Inches(5.2), Inches(5), Inches(0.45),
    size=13, bold=True, color=RGBColor(0x69,0xC3,0xD2))
steps_txt = "1.  Choisir une formule   →   2.  Envoyer les infos du club & photos athlètes   →   3.  Site en ligne sous 10 jours   →   4.  Premier don reçu 🎉"
txt(sl, steps_txt, Inches(1.8), Inches(5.65), Inches(10), Inches(1.1),
    size=13, color=WHITE)

# ── Save ──────────────────────────────────────────────────────────
out = "/Users/klaassabbe/Sites/supsport-club/public/SupSport_Club_Proposition_Partenariat.pptx"
prs.save(out)
print(f"✅  Fichier créé : {out}")
